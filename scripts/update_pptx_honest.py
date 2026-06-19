"""Update deepNoC defense deck: leaky 0.943 -> honest 0.927 NoCNet-v2 numbers."""
from pptx import Presentation
from copy import deepcopy

SRC = "/mnt/c/Users/Administrator/Desktop/deepNoC_NoCNet-v2_Defense.pptx"

# (slide_idx, shape_name, old_substring, new_substring)
EDITS = [
    # headline accuracy 94.3 -> 92.7
    (0,  "Text 6",  "94.3", "92.7"),
    (1,  "Text 21", "94.3", "92.7"),
    (16, "Text 3",  "94.3", "92.7"),
    (22, "Text 1",  "94.3", "92.7"),
    (24, "Text 27", "94.3", "92.7"),
    # slide 9: macro-F1 gain from bias (honest 0.59 -> 0.65)
    (9,  "Text 7",  "0.59 → 0.70", "0.59 → 0.65"),
    # slide 22 progression (hybrid / FT / FT+TTA / FT+TTA+bias)
    (22, "Text 6",  "0.876", "0.856"),
    (22, "Text 11", "0.891", "0.880"),
    (22, "Text 16", "0.891", "0.881"),
    (22, "Text 21", "0.943", "0.927"),
    (22, "Text 23", "+6.7 điểm", "+4.7 điểm"),
    # slide 22 per-class table  NoC1
    (22, "Text 35", "0.994", "1.000"),
    (22, "Text 36", "0.993", "0.991"),
    (22, "Text 37", "0.994", "1.000"),
    (22, "Text 38", "0.993", "0.996"),
    # NoC2
    (22, "Text 44", "0.521", "0.333"),
    (22, "Text 45", "0.658", "0.500"),
    (22, "Text 46", "0.521", "0.333"),
    (22, "Text 47", "0.581", "0.400"),
    (22, "Text 50", "△", "✗"),   # status △ -> ✗
    # NoC3
    (22, "Text 53", "0.844", "0.672"),
    (22, "Text 54", "0.659", "0.566"),
    (22, "Text 55", "0.844", "0.672"),
    (22, "Text 56", "0.740", "0.614"),
    (22, "Text 59", "✓", "△"),   # status ✓ -> △
    # NoC4
    (22, "Text 62", "0.143", "0.214"),
    (22, "Text 63", "0.400", "0.600"),
    (22, "Text 64", "0.143", "0.214"),
    (22, "Text 65", "0.211", "0.316"),
    # NoC5
    (22, "Text 71", "0.962", "0.972"),
    (22, "Text 72", "0.962", "0.912"),
    (22, "Text 73", "0.962", "0.972"),
    (22, "Text 74", "0.962", "0.941"),
    # Overall
    (22, "Text 80", "0.943", "0.927"),
    # slide 23 takeaways
    (23, "Text 6",  "+27.5 điểm", "+26.7 điểm"),
    (23, "Text 9",  "+27.5 điểm", "+25.9 điểm"),
    (23, "Text 12", "+5.2 điểm",  "+4.7 điểm"),
    (23, "Text 16", "acc chỉ 0.143",   "acc chỉ 0.214"),
    # slide 24 conclusion
    (24, "Text 24", "+5.2 điểm",  "+4.7 điểm"),
    (24, "Text 32", "acc 0.143",            "acc 0.214"),
]

# slide 23 comparison chart: hybrid/FT/FT+TTA+Bias
CHART_NEW = [0.66, 0.668, 0.856, 0.880, 0.927]


def replace_in_shape(shape, old, new):
    tf = shape.text_frame
    # try run-level (preserves formatting)
    for para in tf.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                return True
    # fallback: substring spans runs -> rewrite paragraph keeping run0 font
    for para in tf.paragraphs:
        full = "".join(r.text for r in para.runs)
        if old in full:
            if para.runs:
                para.runs[0].text = full.replace(old, new)
                for r in para.runs[1:]:
                    r.text = ""
            return True
    return False


prs = Presentation(SRC)
slides = list(prs.slides)
by_name = []
for s in slides:
    d = {}
    for sh in s.shapes:
        d[sh.name] = sh
    by_name.append(d)

misses = []
for si, name, old, new in EDITS:
    sh = by_name[si].get(name)
    if sh is None or not sh.has_text_frame:
        misses.append((si, name, "NO SHAPE/TF", old))
        continue
    if not replace_in_shape(sh, old, new):
        misses.append((si, name, "OLD NOT FOUND", old))
    else:
        print(f"OK  s{si} {name}: {old!r} -> {new!r}")

# chart
for sh in slides[23].shapes:
    if sh.has_chart:
        ch = sh.chart
        plot = ch.plots[0]
        cats = list(plot.categories)
        from pptx.chart.data import CategoryChartData
        cd = CategoryChartData()
        cd.categories = cats
        cd.add_series(plot.series[0].name, CHART_NEW)
        ch.replace_data(cd)
        print(f"OK  s23 chart {sh.name}: {CHART_NEW}")

print("\nMISSES:", misses if misses else "none")
prs.save(SRC)
print("SAVED", SRC)
